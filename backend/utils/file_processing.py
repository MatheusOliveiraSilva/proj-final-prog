import os
import mimetypes
import uuid
from datetime import datetime
from typing import Dict, Any, Optional, Tuple
import logging
from pathlib import Path

# Configure logging
logger = logging.getLogger(__name__)

def detect_file_type(filename: str, content: bytes) -> Tuple[str, str]:
    """
    Detect file type from filename and content.
    
    Args:
        filename (str): Original filename
        content (bytes): File content
        
    Returns:
        Tuple[str, str]: (mime_type, file_extension)
    """
    # Get mime type from filename
    mime_type, _ = mimetypes.guess_type(filename)
    
    # Get file extension
    file_extension = Path(filename).suffix.lower()
    
    # Fallback to content-based detection for common types
    if not mime_type:
        if content.startswith(b'%PDF'):
            mime_type = 'application/pdf'
            file_extension = '.pdf'
        elif content.startswith(b'PK\x03\x04'):  # ZIP-based formats (DOCX, XLSX, etc.)
            if file_extension in ['.docx', '.doc']:
                mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            elif file_extension in ['.xlsx', '.xls']:
                mime_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            elif file_extension in ['.pptx', '.ppt']:
                mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        elif content.startswith(b'\xd0\xcf\x11\xe0'):  # Old Office formats
            if file_extension == '.doc':
                mime_type = 'application/msword'
            elif file_extension == '.xls':
                mime_type = 'application/vnd.ms-excel'
            elif file_extension == '.ppt':
                mime_type = 'application/vnd.ms-powerpoint'
        else:
            # Try to decode as text
            try:
                content.decode('utf-8')
                mime_type = 'text/plain'
                if not file_extension:
                    file_extension = '.txt'
            except UnicodeDecodeError:
                mime_type = 'application/octet-stream'
    
    return mime_type or 'application/octet-stream', file_extension


def extract_text_content(content: bytes, mime_type: str, filename: str) -> str:
    """
    Extract text content from different file types.
    
    Args:
        content (bytes): File content
        mime_type (str): MIME type
        filename (str): Original filename
        
    Returns:
        str: Extracted text content
    """
    try:
        if mime_type == 'text/plain':
            return content.decode('utf-8')
        
        elif mime_type == 'application/pdf':
            return extract_pdf_content(content)
        
        elif 'word' in mime_type or filename.endswith(('.doc', '.docx')):
            return extract_word_content(content, filename)
        
        elif 'spreadsheet' in mime_type or filename.endswith(('.xls', '.xlsx')):
            return extract_excel_content(content, filename)
        
        elif 'presentation' in mime_type or filename.endswith(('.ppt', '.pptx')):
            return extract_powerpoint_content(content, filename)
        
        elif mime_type.startswith('text/'):
            # Try to decode as text with different encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    return content.decode(encoding)
                except UnicodeDecodeError:
                    continue
            return content.decode('utf-8', errors='ignore')
        
        else:
            # Fallback: try to decode as text
            try:
                return content.decode('utf-8')
            except UnicodeDecodeError:
                logger.warning(f"Could not extract text from {mime_type} file: {filename}")
                return f"[Binary file: {filename}]"
                
    except Exception as e:
        logger.error(f"Error extracting content from {filename}: {str(e)}")
        return f"[Error extracting content from {filename}]"


def extract_pdf_content(content: bytes) -> str:
    """Extract text from PDF content."""
    try:
        import PyPDF2
        import io
        
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
        text_content = []
        
        for page in pdf_reader.pages:
            text_content.append(page.extract_text())
        
        return '\n'.join(text_content)
        
    except ImportError:
        logger.warning("PyPDF2 not installed. Install with: pip install PyPDF2")
        return "[PDF content - PyPDF2 not available]"
    except Exception as e:
        logger.error(f"Error extracting PDF content: {str(e)}")
        return "[Error extracting PDF content]"


def extract_word_content(content: bytes, filename: str) -> str:
    """Extract text from Word documents."""
    try:
        import docx
        import io
        
        if filename.endswith('.docx'):
            doc = docx.Document(io.BytesIO(content))
            text_content = []
            
            for paragraph in doc.paragraphs:
                text_content.append(paragraph.text)
            
            return '\n'.join(text_content)
        else:
            # For .doc files, we'd need python-docx2txt or other library
            logger.warning("Legacy .doc format not fully supported. Please convert to .docx")
            return "[Legacy .doc file - limited support]"
            
    except ImportError:
        logger.warning("python-docx not installed. Install with: pip install python-docx")
        return "[Word content - python-docx not available]"
    except Exception as e:
        logger.error(f"Error extracting Word content: {str(e)}")
        return "[Error extracting Word content]"


def extract_excel_content(content: bytes, filename: str) -> str:
    """Extract text from Excel spreadsheets."""
    try:
        import pandas as pd
        import io
        
        # Read Excel file
        if filename.endswith('.xlsx'):
            xl_file = pd.ExcelFile(io.BytesIO(content))
        else:
            xl_file = pd.ExcelFile(io.BytesIO(content), engine='xlrd')
        
        text_content = []
        
        for sheet_name in xl_file.sheet_names:
            df = xl_file.parse(sheet_name)
            text_content.append(f"Sheet: {sheet_name}")
            text_content.append(df.to_string(index=False))
            text_content.append("\n")
        
        return '\n'.join(text_content)
        
    except ImportError:
        logger.warning("pandas/xlrd not installed. Install with: pip install pandas xlrd openpyxl")
        return "[Excel content - pandas not available]"
    except Exception as e:
        logger.error(f"Error extracting Excel content: {str(e)}")
        return "[Error extracting Excel content]"


def extract_powerpoint_content(content: bytes, filename: str) -> str:
    """Extract text from PowerPoint presentations."""
    try:
        from pptx import Presentation
        import io
        
        prs = Presentation(io.BytesIO(content))
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_content.append(f"Slide {slide_num}:")
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
            
            text_content.append("\n")
        
        return '\n'.join(text_content)
        
    except ImportError:
        logger.warning("python-pptx not installed. Install with: pip install python-pptx")
        return "[PowerPoint content - python-pptx not available]"
    except Exception as e:
        logger.error(f"Error extracting PowerPoint content: {str(e)}")
        return "[Error extracting PowerPoint content]"


def generate_document_metadata(
    filename: str, 
    content: str, 
    mime_type: str, 
    file_size: int,
    namespace: str = ""
) -> Dict[str, Any]:
    """
    Generate document metadata automatically from file information.
    
    Args:
        filename (str): Original filename
        content (str): Extracted text content
        mime_type (str): MIME type
        file_size (int): File size in bytes
        namespace (str): Namespace/thread ID
        
    Returns:
        Dict[str, Any]: Generated metadata
    """
    # Generate document ID
    doc_id = str(uuid.uuid4())
    
    # Extract title from filename (remove extension)
    title = Path(filename).stem
    
    # Generate creation timestamp
    created_at = datetime.utcnow().isoformat()
    
    # Determine document type from mime type
    document_type = "document"
    if "pdf" in mime_type:
        document_type = "pdf"
    elif "word" in mime_type:
        document_type = "word_document"
    elif "spreadsheet" in mime_type:
        document_type = "spreadsheet"
    elif "presentation" in mime_type:
        document_type = "presentation"
    elif "text" in mime_type:
        document_type = "text_file"
    
    # Generate basic tags from filename and content
    tags = []
    file_ext = Path(filename).suffix.lower().replace('.', '')
    if file_ext:
        tags.append(file_ext)
    
    # Add content-based tags (simple keyword extraction)
    content_lower = content.lower()
    keyword_tags = []
    if any(word in content_lower for word in ['machine learning', 'ml', 'ai', 'artificial intelligence']):
        keyword_tags.append('machine_learning')
    if any(word in content_lower for word in ['python', 'programming', 'code']):
        keyword_tags.append('programming')
    if any(word in content_lower for word in ['data', 'analysis', 'statistics']):
        keyword_tags.append('data_analysis')
    
    tags.extend(keyword_tags)
    
    # Calculate content statistics
    word_count = len(content.split())
    char_count = len(content)
    
    metadata = {
        'id': doc_id,
        'title': title,
        'content': content,
        'source': 'file_upload',
        'filename': filename,
        'mime_type': mime_type,
        'document_type': document_type,
        'file_size': file_size,
        'created_at': created_at,
        'uploaded_at': created_at,
        'namespace': namespace,
        'tags': tags,
        'word_count': word_count,
        'char_count': char_count,
        'author': None,  # Could be extracted from file metadata in the future
    }
    
    logger.info(f"Generated metadata for {filename}: {word_count} words, {len(tags)} tags")
    
    return metadata


def validate_file_upload(filename: str, content: bytes, max_size_mb: int = 50) -> Tuple[bool, str]:
    """
    Validate file upload constraints.
    
    Args:
        filename (str): Original filename
        content (bytes): File content
        max_size_mb (int): Maximum file size in MB
        
    Returns:
        Tuple[bool, str]: (is_valid, error_message)
    """
    # Check file size
    file_size_mb = len(content) / (1024 * 1024)
    if file_size_mb > max_size_mb:
        return False, f"File size ({file_size_mb:.1f}MB) exceeds maximum allowed size ({max_size_mb}MB)"
    
    # Check filename
    if not filename or filename.strip() == "":
        return False, "Filename cannot be empty"
    
    # Check for dangerous file extensions
    dangerous_extensions = ['.exe', '.bat', '.cmd', '.com', '.scr', '.vbs']
    file_ext = Path(filename).suffix.lower()
    if file_ext in dangerous_extensions:
        return False, f"File type {file_ext} is not allowed for security reasons"
    
    # Check if content is empty
    if len(content) == 0:
        return False, "File is empty"
    
    return True, "" 